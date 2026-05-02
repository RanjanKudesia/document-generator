"""PPTX document generation pipeline."""
import base64
import logging
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

# pylint: disable=broad-exception-caught

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.document_generation_schema import (
    DocumentGenerationRequest,
    ExtractedData,
    ExtractedPptData,
    ExtractedXmlData,
    ParagraphBlock,
    TableBlock,
)


class PptGenerationPipeline:
    """Generate PPTX documents from blocks or extracted payloads.

    Notes:
    - Uses a practical reconstruction strategy for extracted payloads.
    - JSON extracted_data: renders ordered paragraphs and tables into content slides.
    - XML extracted_data: renders parsed_body paragraphs/tables into content slides.
    """

    def __init__(self) -> None:
        self.logger = logging.getLogger(__name__)

    def run(self, payload: DocumentGenerationRequest, file_name: str) -> bytes:
        """Build a PPTX from the payload and return raw bytes."""
        _ = file_name
        data_type = type(
            payload.extracted_data).__name__ if payload.extracted_data else "blocks"
        self.logger.info(
            "ppt_pipeline_start file=%s data_type=%s", file_name, data_type)
        # Highest-fidelity path: reconstruct pptx package directly from extracted
        # OpenXML parts when available (preserves masters/themes/backgrounds).
        if isinstance(payload.extracted_data, ExtractedPptData):
            rebuilt = self._try_rebuild_from_package_dump(
                payload.extracted_data)
            if rebuilt is not None:
                self.logger.info(
                    "ppt_pipeline_complete file=%s strategy=package_dump size_bytes=%d",
                    file_name, len(rebuilt)
                )
                return rebuilt

        prs = Presentation()

        if payload.extracted_data is not None:
            if isinstance(payload.extracted_data, ExtractedPptData):
                self._from_ppt_extracted(
                    prs, payload.extracted_data, payload.title)
            elif isinstance(payload.extracted_data, ExtractedXmlData):
                self._from_xml(prs, payload.extracted_data, payload.title)
            else:
                self._from_json(prs, payload.extracted_data, payload.title)
        else:
            self._from_blocks(prs, payload)

        with BytesIO() as output:
            prs.save(output)
            result = output.getvalue()
        self.logger.info(
            "ppt_pipeline_complete file=%s strategy=generated size_bytes=%d",
            file_name, len(result)
        )
        return result

    def _try_rebuild_from_package_dump(  # pylint: disable=too-many-branches  # NOSONAR
        self, data: ExtractedPptData
    ) -> bytes | None:
        if not data.parts:
            self.logger.debug("ppt_package_dump_skip reason=no_parts")
            return None

        xml_parts = {
            (p.path or "").lstrip("/"): (p.xml or "")
            for p in data.parts
            if (p.path or "")
        }

        # Package-level minimum for valid pptx reconstruction.
        required = ["[Content_Types].xml",
                    "_rels/.rels", "ppt/presentation.xml"]
        if not all(k in xml_parts for k in required):
            self.logger.debug(
                "ppt_package_dump_skip reason=missing_required_parts xml_parts=%d",
                len(xml_parts)
            )
            return None

        media_bytes_by_path: dict[str, bytes] = {}
        for media in data.media:
            b64 = getattr(media, "base64_data", None) or getattr(
                media, "base64", None)
            path = (getattr(media, "local_file_path", None) or "").lstrip("/")
            if not b64 or not path:
                continue
            try:
                media_bytes_by_path[path] = base64.b64decode(
                    b64, validate=True)
            except (TypeError, ValueError):
                self.logger.debug("ppt_media_decode_failed path=%s", path)
                continue

        # XML extraction mode stores picture blobs under parsed_slides[*].shapes[*]
        # rather than top-level media[]. Rehydrate those too.
        for slide in (data.parsed_slides or []):
            for shape in (slide.shapes or []):
                if shape.get("kind") != "picture":
                    continue
                path = (shape.get("target_path") or "").lstrip("/")
                b64 = shape.get("base64")
                if not path or not b64:
                    continue
                # Keep first occurrence for determinism.
                if path in media_bytes_by_path:
                    continue
                try:
                    media_bytes_by_path[path] = base64.b64decode(
                        b64, validate=True)
                except (TypeError, ValueError):
                    self.logger.debug(
                        "ppt_shape_media_decode_failed path=%s", path)
                    continue

        binary_bytes_by_path: dict[str, bytes] = {}
        for item in (data.binary_parts or []):
            path = (item.get("path") or "").lstrip(
                "/") if isinstance(item, dict) else ""
            b64 = item.get("base64") if isinstance(item, dict) else None
            if not path or not b64:
                continue
            try:
                binary_bytes_by_path[path] = base64.b64decode(
                    b64, validate=True)
            except (TypeError, ValueError):
                self.logger.debug(
                    "ppt_binary_part_decode_failed path=%s", path)
                continue

        output = BytesIO()
        with ZipFile(output, "w", compression=ZIP_DEFLATED) as archive:
            for path, xml in xml_parts.items():
                archive.writestr(path, xml)

            for path, blob in binary_bytes_by_path.items():
                # Avoid clobbering explicit media write below.
                if path in media_bytes_by_path:
                    continue
                archive.writestr(path, blob)

            for path, blob in media_bytes_by_path.items():
                archive.writestr(path, blob)

        self.logger.debug(
            "ppt_package_dump_assembled xml_parts=%d media=%d binary=%d",
            len(xml_parts), len(media_bytes_by_path), len(binary_bytes_by_path)
        )
        return output.getvalue()

    def _from_ppt_extracted(  # pylint: disable=too-many-branches,too-many-statements  # NOSONAR
        self, prs: Presentation, data: ExtractedPptData, title: str | None
    ) -> None:
        # JSON-adapter path: slide-wise rehydration using slide indices.
        if data.slides:
            self.logger.info(
                "ppt_from_ppt_extracted path=json_slides slide_count=%d",
                len(data.slides)
            )
            paragraphs_by_idx = {p.index: p for p in data.paragraphs}
            tables_by_idx = {t.index: t for t in data.tables}
            media_by_idx = dict(enumerate(data.media))

            for slide in sorted(
                data.slides, key=lambda s: (
                    s.index if s.index is not None else 10**9)
            ):
                slide_title = (
                    slide.title
                    or f"Slide {(slide.slide_number or ((slide.index or 0) + 1))}"
                ).strip()
                slide_lines: list[str] = []

                for p_idx in slide.paragraph_indices:
                    para = paragraphs_by_idx.get(p_idx)
                    if para is not None and (para.text or "").strip():
                        slide_lines.append((para.text or "").strip())

                if not slide_lines and (slide.text or "").strip():
                    slide_lines = [line.strip() for line in (
                        slide.text or "").splitlines() if line.strip()]

                # Avoid duplicating title text in body content.
                slide_lines = self._clean_body_lines(
                    slide_lines, slide_title)

                slide_tables: list[list[list[str]]] = []
                for t_idx in slide.table_indices:
                    tbl = tables_by_idx.get(t_idx)
                    if tbl is None:
                        continue
                    rows = [[(cell.text or "") for cell in row.cells]
                            for row in tbl.rows]
                    if rows:
                        slide_tables.append(rows)

                slide_media = []
                for m_idx in slide.media_indices:
                    m = media_by_idx.get(m_idx)
                    if m is not None and not self._is_placeholder_media(m):
                        slide_media.append(m)

                # Keep one generated slide per source slide to preserve ordering/fidelity.
                title_color_rgb = self._pick_title_color_from_paragraphs(
                    slide_title, slide.paragraph_indices, paragraphs_by_idx)
                out_slide = self._add_composite_slide(
                    prs,
                    title=slide_title or "Slide",
                    lines=slide_lines,
                    tables=slide_tables,
                    media_items=slide_media,
                    title_color_rgb=title_color_rgb,
                )
                self._pad_to_shape_count(out_slide, slide.shape_count)
                self._write_speaker_notes(
                    out_slide, getattr(slide, "notes_text", None))

            return

        # XML-pipeline path: use parsed_slides directly.
        if data.parsed_slides:
            self.logger.info(
                "ppt_from_ppt_extracted path=parsed_slides slide_count=%d",
                len(data.parsed_slides)
            )
            for slide in sorted(
                data.parsed_slides, key=lambda s: (
                    s.index if s.index is not None else 10**9)
            ):
                if slide.parse_error:
                    self.logger.warning(
                        "ppt_slide_skipped reason=parse_error slide_index=%s error=%s",
                        slide.index, slide.parse_error
                    )
                    continue

                slide_title = (
                    slide.title or f"Slide {((slide.index or 0) + 1)}").strip()
                slide_lines: list[str] = []

                if (slide.text or "").strip():
                    slide_lines = [line.strip() for line in (
                        slide.text or "").splitlines() if line.strip()]

                    slide_lines = self._clean_body_lines(
                        slide_lines, slide_title)

                if slide.notes and isinstance(slide.notes, dict):
                    notes_text = (slide.notes.get("text") or "").strip()
                    if notes_text:
                        slide_lines.append(f"Notes: {notes_text}")

                slide_tables: list[list[list[str]]] = []
                for shape in (slide.shapes or []):
                    if (
                        shape.get("kind") != "graphic_frame"
                        or shape.get("graphic_type") != "table"
                    ):
                        continue
                    table_payload = shape.get("table") or {}
                    rows = []
                    for row in table_payload.get("rows", []) or []:
                        rows.append([(cell.get("text") or "")
                                    for cell in (row.get("cells") or [])])
                    if rows:
                        slide_tables.append(rows)

                out_slide = self._add_composite_slide(
                    prs,
                    title=slide_title or "Slide",
                    lines=slide_lines,
                    tables=slide_tables,
                    media_items=[],
                    title_color_rgb=self._pick_title_color_from_shapes(
                        slide.shapes),
                )
                self._pad_to_shape_count(out_slide, slide.shape_count)
                notes_text = (slide.notes.get("text") or "").strip(
                ) if isinstance(slide.notes, dict) else ""
                self._write_speaker_notes(out_slide, notes_text)

            return

        # Fallback for malformed payloads: route to generic handlers if possible.
        if data.paragraphs or data.tables:
            self._from_json(prs, ExtractedData(
                document_order=data.document_order,
                styles=data.styles,
                paragraphs=data.paragraphs,
                tables=data.tables,
                media=data.media,
            ), title)

    def _from_blocks(self, prs: Presentation, payload: DocumentGenerationRequest) -> None:
        if payload.title:
            self._add_title_slide(prs, payload.title, "")

        for block in payload.blocks:
            if isinstance(block, ParagraphBlock):
                title = f"Heading {block.heading_level}" if block.heading_level else "Paragraph"
                self._add_text_slide(prs, title=title, lines=[block.text])
            elif isinstance(block, TableBlock):
                self._add_table_slide(prs, title="Table", rows=block.rows)

    def _from_json(  # NOSONAR
        self, prs: Presentation, data: ExtractedData, title: str | None
    ) -> None:
        self.logger.info(
            "ppt_from_json paragraphs=%d tables=%d",
            len(data.paragraphs), len(data.tables)
        )
        if title:
            self._add_title_slide(prs, title, "")

        para_by_idx = {p.index: p for p in data.paragraphs}
        table_by_idx = {t.index: t for t in data.tables}

        if data.document_order:
            for item in data.document_order:
                if item.type == "paragraph":
                    p = para_by_idx.get(item.index)
                    if p is None:
                        continue
                    heading = p.style if p.style else "Paragraph"
                    text = p.text or ""
                    if text.strip():
                        self._add_text_slide(prs, title=heading, lines=[text])
                elif item.type == "table":
                    t = table_by_idx.get(item.index)
                    if t is None:
                        continue
                    rows = [[(cell.text or "") for cell in row.cells]
                            for row in t.rows]
                    self._add_table_slide(prs, title="Table", rows=rows)
            return

        for p in sorted(data.paragraphs, key=lambda x: x.index):
            text = p.text or ""
            if not text.strip():
                continue
            self._add_text_slide(prs, title=(
                p.style or "Paragraph"), lines=[text])

        for t in sorted(data.tables, key=lambda x: x.index):
            rows = [[(cell.text or "") for cell in row.cells]
                    for row in t.rows]
            self._add_table_slide(prs, title="Table", rows=rows)

    def _from_xml(  # NOSONAR
        self, prs: Presentation, data: ExtractedXmlData, title: str | None
    ) -> None:
        self.logger.info(
            "ppt_from_xml body_items=%d",
            len(data.parsed_body)
        )
        if title:
            self._add_title_slide(prs, title, "")

        for item in data.parsed_body:
            if item.type == "paragraph" and item.paragraph is not None:
                p = item.paragraph
                text = p.text or ""
                if not text and p.runs:
                    text = "".join((r.text or "") for r in p.runs)
                if text.strip():
                    heading = p.style_id or "Paragraph"
                    self._add_text_slide(prs, title=heading, lines=[text])
            elif item.type == "table" and item.table is not None:
                rows = [[(cell.text or "") for cell in row.cells]
                        for row in item.table.rows]
                self._add_table_slide(prs, title="Table", rows=rows)

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def _add_text_slide(self, prs: Presentation, title: str, lines: list[str]) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        first = True
        for line in lines:
            text = (line or "").strip()
            if not text:
                continue

            if first:
                p = body.paragraphs[0]
                first = False
            else:
                p = body.add_paragraph()

            p.text = text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.size = Pt(18)

    def _add_table_slide(self, prs: Presentation, title: str, rows: list[list[str]]) -> None:
        layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

        if not rows:
            return

        col_count = max((len(r) for r in rows), default=0)
        if col_count == 0:
            return

        norm_rows = [r + [""] * (col_count - len(r)) for r in rows]

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.0)

        table_shape = slide.shapes.add_table(
            rows=len(norm_rows),
            cols=col_count,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        table = table_shape.table

        for r_i, row in enumerate(norm_rows):
            for c_i, value in enumerate(row):
                table.cell(r_i, c_i).text = value or ""

    def _add_composite_slide(  # NOSONAR
        self,
        prs: Presentation,
        title: str,
        lines: list[str],
        tables: list[list[list[str]]],
        media_items: list,
        title_color_rgb: str | None = None,
    ):
        """Create one slide that can contain title, text, media, and at most one table.

        This preserves source slide count and avoids shifting content across slides.
        """
        layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
            self._apply_title_color(slide.shapes.title, title_color_rgb)

        has_table = bool(tables)
        has_media = bool(media_items)

        # Text region
        text_lines = [ln for ln in lines if (ln or "").strip()]

        # Only add a body textbox when there is real body text.
        if text_lines:
            left = Inches(0.5)
            top = Inches(1.2)
            text_width = Inches(6.2 if has_media else 9.0)
            text_height = Inches(2.4 if has_table else 5.4)

            textbox = slide.shapes.add_textbox(
                left, top, text_width, text_height)
            tf = textbox.text_frame
            tf.clear()

            first = True
            for line in text_lines:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = (line or "").strip()
                p.level = 0
                p.alignment = PP_ALIGN.LEFT
                if p.runs:
                    p.runs[0].font.size = Pt(16)

        # Media region (right column)
        if has_media:
            media_left = Inches(7.0)
            media_top = Inches(1.2)
            max_w = Inches(2.4)
            max_h = Inches(1.6)

            for media in media_items[:3]:
                added = self._add_media_image(
                    slide,
                    media,
                    media_left,
                    media_top,
                    max_w,
                    max_h,
                )
                if added:
                    media_top += Inches(1.8)

        # Table region (bottom)
        if has_table:
            self._add_table_shape(
                slide,
                rows=tables[0],
                left=Inches(0.5),
                top=Inches(3.9),
                width=Inches(9.0),
                height=Inches(2.8),
            )

        return slide

    def _apply_title_color(self, title_shape, color_rgb: str | None) -> None:
        color = self._normalize_hex_rgb(color_rgb)
        if not color:
            return
        try:
            tf = title_shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor.from_string(color)
        except Exception:  # NOSONAR
            return

    def _normalize_hex_rgb(self, value: str | None) -> str | None:
        if not value:
            return None
        v = value.strip()
        if v.lower().startswith("scheme:"):
            return None
        if v.startswith("#"):
            v = v[1:]
        if len(v) != 6:
            return None
        try:
            int(v, 16)
        except ValueError:
            return None
        return v.upper()

    def _pick_title_color_from_paragraphs(  # NOSONAR
        self,
        slide_title: str,
        paragraph_indices: list[int],
        paragraphs_by_idx: dict[int, object],
    ) -> str | None:
        title_norm = (slide_title or "").strip().lower()
        for idx in paragraph_indices:
            para = paragraphs_by_idx.get(idx)
            if para is None:
                continue
            text = (getattr(para, "text", "") or "").strip()
            if not text:
                continue
            if title_norm and text.lower() != title_norm:
                continue
            for run in getattr(para, "runs", []) or []:
                color = getattr(run, "color_rgb", None)
                if self._normalize_hex_rgb(color):
                    return color
        # fallback: first explicit run color on slide
        for idx in paragraph_indices:
            para = paragraphs_by_idx.get(idx)
            if para is None:
                continue
            for run in getattr(para, "runs", []) or []:
                color = getattr(run, "color_rgb", None)
                if self._normalize_hex_rgb(color):
                    return color
        return None

    def _pick_title_color_from_shapes(self, shapes: list[dict]) -> str | None:
        for shape in shapes or []:
            if not shape.get("is_title"):
                continue
            for para in shape.get("paragraphs", []) or []:
                for run in para.get("runs", []) or []:
                    color = run.get("color_rgb")
                    if self._normalize_hex_rgb(color):
                        return color
        return None

    def _write_speaker_notes(self, slide, notes_text: str | None) -> None:
        """Write speaker notes to a slide's notes text frame."""
        if not notes_text or not notes_text.strip():
            return
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text.strip()
        except Exception:  # NOSONAR
            pass

    def _pad_to_shape_count(self, slide, expected_count: int | None) -> None:
        if expected_count is None:
            return
        missing = expected_count - len(slide.shapes)
        if missing <= 0:
            return

        # Keep structure closer to source deck without adding visible content.
        for _ in range(missing):
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9.8),
                Inches(7.2),
                Inches(0.01),
                Inches(0.01),
            )
            shp.fill.background()
            shp.line.fill.background()

    def _add_table_shape(self, slide, rows: list[list[str]], left, top, width, height) -> None:
        if not rows:
            return
        col_count = max((len(r) for r in rows), default=0)
        if col_count == 0:
            return
        norm_rows = [r + [""] * (col_count - len(r)) for r in rows]
        table_shape = slide.shapes.add_table(
            rows=len(norm_rows),
            cols=col_count,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        table = table_shape.table
        for r_i, row in enumerate(norm_rows):
            for c_i, value in enumerate(row):
                table.cell(r_i, c_i).text = value or ""

    def _add_media_image(self, slide, media, left, top, max_w, max_h) -> bool:
        """Insert media image if possible. Returns True when picture was added."""
        img_blob = None
        # Support either base64_data or base64 field from extractor payloads.
        b64 = getattr(media, "base64_data", None) or getattr(
            media, "base64", None)
        if b64:
            try:
                img_blob = base64.b64decode(b64, validate=True)
            except (TypeError, ValueError):
                img_blob = None

        if not img_blob:
            return False

        try:
            stream = BytesIO(img_blob)
            slide.shapes.add_picture(
                stream, left, top, width=max_w, height=max_h)
            return True
        except Exception:  # NOSONAR
            return False

    def _clean_body_lines(self, lines: list[str], title: str) -> list[str]:
        """Remove noise/duplicates from body text while preserving order."""
        out: list[str] = []
        title_norm = (title or "").strip().lower()
        seen: set[str] = set()
        for ln in lines:
            t = (ln or "").strip()
            if not t:
                continue
            t_norm = t.lower()
            if title_norm and t_norm == title_norm:
                continue
            if t_norm in seen:
                continue
            seen.add(t_norm)
            out.append(t)
        return out

    def _is_placeholder_media(self, media) -> bool:
        src = getattr(media, "source", None)
        name = ""
        if isinstance(src, dict):
            name = (src.get("name") or "").lower()
        elif src is not None:
            # Pydantic models expose source metadata as attributes.
            name = (getattr(src, "name", "") or "").lower()

        if "placeholder" in name:
            return True
        return False
